import os
import subprocess
from copy import deepcopy
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

# Pygments for code highlighting
from pygments import highlight
from pygments.lexers import get_lexer_by_name
from pygments.formatters import ImageFormatter

import logging

logger = logging.getLogger("deassignment.ppt")
logger.setLevel(logging.DEBUG)
if not logger.handlers:
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter("[%(asctime)s] [PPT] %(levelname)s: %(message)s", datefmt="%H:%M:%S"))
    logger.addHandler(handler)


class PPTEngine:
    """Handles all PPTX generation and manipulation."""

    # Preferred monospace fonts in priority order
    _MONO_FONT_CANDIDATES = [
        "DejaVu Sans Mono",
        "Liberation Mono",
        "Noto Sans Mono",
        "Courier New",
        "Cousine",
        "monospace",
    ]

    # macOS-specific font candidates (Menlo is always available on macOS)
    _MACOS_MONO_FONT_CANDIDATES = [
        "DejaVu Sans Mono",
        "Menlo",
        "Monaco",
        "Courier New",
        "Courier",
    ]

    # Common macOS font directories
    _MACOS_FONT_DIRS = [
        "/System/Library/Fonts",
        "/Library/Fonts",
        os.path.expanduser("~/Library/Fonts"),
    ]

    _resolved_font = None  # Class-level cache so we only resolve once

    def __init__(self, template_path: str, config: dict):
        self.prs = Presentation(template_path)
        self.config = config
        
        # Determine the template slide (user requested last slide)
        if len(self.prs.slides) == 0:
            raise ValueError("Template presentation has no slides.")
        self.template_slide = self.prs.slides[-1]

        # Resolve the monospace font once at startup
        if PPTEngine._resolved_font is None:
            PPTEngine._resolved_font = self._detect_mono_font()
            logger.info(f"Resolved monospace font: '{PPTEngine._resolved_font}'")
        
    @classmethod
    def _detect_mono_font(cls) -> str:
        """Find the first available monospace font on the system.

        On macOS: checks font directories directly and uses macOS-specific candidates.
        On Linux: uses fc-list (fontconfig) with fallback to PIL font loading.
        """
        import sys

        if sys.platform == 'darwin':
            return cls._detect_mono_font_macos()

        # Linux: use fc-list
        for font_name in cls._MONO_FONT_CANDIDATES:
            try:
                result = subprocess.run(
                    ["fc-list", font_name],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=3
                )
                output = result.stdout.decode().strip()
                if output:
                    logger.debug(f"Font '{font_name}' found on system")
                    return font_name
                else:
                    logger.debug(f"Font '{font_name}' NOT found on system")
            except (FileNotFoundError, subprocess.TimeoutExpired):
                # fc-list not available, try PIL/Pillow font loading
                try:
                    from PIL import ImageFont
                    ImageFont.truetype(font_name, 16)
                    logger.debug(f"Font '{font_name}' loadable via PIL")
                    return font_name
                except (IOError, OSError):
                    logger.debug(f"Font '{font_name}' not loadable via PIL either")
                    continue

        # Ultimate fallback — Pygments will use its built-in default
        logger.warning(
            "No preferred monospace font found on system! "
            "Pygments will fall back to its built-in default. "
            "Install a monospace font: sudo apt install fonts-dejavu-core "
            "or sudo dnf install dejavu-sans-mono-fonts"
        )
        return ""

    @classmethod
    def _detect_mono_font_macos(cls) -> str:
        """macOS-specific font detection.

        Checks system font directories and uses macOS-native font candidates.
        Menlo is always present on macOS 10.6+ as a built-in monospace font.
        """
        for font_name in cls._MACOS_MONO_FONT_CANDIDATES:
            # Check if font files exist in known macOS font directories
            for font_dir in cls._MACOS_FONT_DIRS:
                if not os.path.isdir(font_dir):
                    continue
                try:
                    for entry in os.listdir(font_dir):
                        # Match font name in filename (case-insensitive)
                        entry_lower = entry.lower()
                        font_lower = font_name.lower().replace(" ", "")
                        if font_lower in entry_lower and entry_lower.endswith(('.ttf', '.ttc', '.otf')):
                            logger.debug(f"Font '{font_name}' found at {os.path.join(font_dir, entry)}")
                            return font_name
                except OSError:
                    continue

            # Also try loading via PIL/Pillow as fallback
            try:
                from PIL import ImageFont
                ImageFont.truetype(font_name, 16)
                logger.debug(f"Font '{font_name}' loadable via PIL on macOS")
                return font_name
            except (IOError, OSError):
                logger.debug(f"Font '{font_name}' not found on macOS")
                continue

        # Menlo should always be available — return it as ultimate macOS fallback
        logger.warning("Using 'Menlo' as fallback monospace font on macOS")
        return "Menlo"

    def _duplicate_template_slide(self):
        """
        Creates a new slide based on the template slide.
        Properly handles copying both text/shapes and images.
        """
        new_slide = self.prs.slides.add_slide(self.template_slide.slide_layout)
        
        for shape in self.template_slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Re-add image via blob to maintain valid rId relationships
                image_blob = shape.image.blob
                image_stream = io.BytesIO(image_blob)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left, shape.top,
                    shape.width, shape.height
                )
            else:
                # Safe to deepcopy XML for other elements
                new_el = deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                
        return new_slide
        
    def _add_header_to_slide(self, slide, text: str):
        """Adds a title/header to the slide within the top safe zone."""
        left = Inches(self.config.get("safe_zone_left", 0.5))
        top = Inches(self.config.get("safe_zone_top", 1.5))
        width = self.prs.slide_width - Inches(self.config.get("safe_zone_left", 0.5) + self.config.get("safe_zone_right", 0.5))
        height = Inches(1.0)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = self.config.get("font_name", "Arial")
        run.font.size = Pt(self.config.get("font_size_title", 18))
        run.font.bold = True
        return top + height # Return bottom position of header

    def add_question_slide(self, q_num: int, q_statement: str):
        """Adds a slide showing just the question."""
        slide = self._duplicate_template_slide()
        self._add_header_to_slide(slide, f"Question {q_num}:\n{q_statement}")

    def _code_to_image(self, code: str, language: str, output_path: str):
        """Renders code to a syntax-highlighted PNG image using Pygments."""
        lexer = get_lexer_by_name(language)

        # Build formatter kwargs, only set font_name if we found one
        fmt_kwargs = dict(
            font_size=16,
            line_pad=6,
            image_pad=10,
            line_numbers=True,
            style=self.config.get("code_style", "monokai"),
        )
        if PPTEngine._resolved_font:
            fmt_kwargs["font_name"] = PPTEngine._resolved_font

        formatter = ImageFormatter(**fmt_kwargs)

        with open(output_path, "wb") as f:
            highlight(code, lexer, formatter, outfile=f)

    def add_code_slide(self, q_num: int, q_statement: str, code: str, language: str, temp_dir: str, code_format: str = "text"):
        """
        Adds slide(s) containing the code and the question statement.
        Supports rendering as a syntax-highlighted image or pasting as editable highlighted text.
        If code is too long, splits it across multiple slides.
        """
        if code_format == "image":
            # 1. Render entire code to image
            code_img_path = os.path.join(temp_dir, f"code_{q_num}.png")
            self._code_to_image(code, language, code_img_path)
            
            # 2. Setup the first slide with the full question statement
            slide = self._duplicate_template_slide()
            header_text = f"Question {q_num}:\n{q_statement}"
            header_bottom = self._add_header_to_slide(slide, header_text)
            
            left = Inches(self.config.get("safe_zone_left", 0.5))
            top = header_bottom + Inches(0.2)
            
            max_width = self.prs.slide_width - Inches(self.config.get("safe_zone_left", 0.5) + self.config.get("safe_zone_right", 0.5))
            max_height = self.prs.slide_height - top - Inches(self.config.get("safe_zone_bottom", 1.0))
            
            with Image.open(code_img_path) as img:
                img_w, img_h = img.size
                
                # Calculate rendering scale (assume 96 DPI natively = 9525 EMUs per pixel)
                EMUS_PER_PIXEL = 9525
                
                rendered_w = img_w * EMUS_PER_PIXEL
                if rendered_w > max_width:
                    rendered_w = max_width
                    
                scale = rendered_w / (img_w * EMUS_PER_PIXEL)
                rendered_h_full = img_h * EMUS_PER_PIXEL * scale
                
                if rendered_h_full <= max_height:
                    slide.shapes.add_picture(code_img_path, left, top, width=rendered_w)
                else:
                    # Image is too tall, we need to crop and spill over to new slides
                    y_offset = 0
                    slide_idx = 0
                    
                    while y_offset < img_h:
                        if slide_idx == 0:
                            current_slide = slide
                            current_max_h = max_height
                            current_top = top
                        else:
                            current_slide = self._duplicate_template_slide()
                            h_bottom = self._add_header_to_slide(current_slide, f"")
                            current_top = h_bottom + Inches(0.2)
                            current_max_h = self.prs.slide_height - current_top - Inches(self.config.get("safe_zone_bottom", 1.0))
                            
                        # How many pixels can fit in the available height?
                        px_to_fit = int(current_max_h / (EMUS_PER_PIXEL * scale))
                        
                        if px_to_fit <= 50:
                            px_to_fit = 200 # Fallback to prevent infinite loop
                            
                        crop_box = (0, y_offset, img_w, min(y_offset + px_to_fit, img_h))
                        cropped = img.crop(crop_box)
                        chunk_path = os.path.join(temp_dir, f"code_chunk_{q_num}_{slide_idx}.png")
                        cropped.save(chunk_path)
                        
                        current_slide.shapes.add_picture(chunk_path, left, current_top, width=rendered_w)
                        
                        y_offset += px_to_fit
                        slide_idx += 1
        else:
            # 1. Parse code into tokens and group by lines
            from pygments.styles import get_style_by_name
            
            lexer = get_lexer_by_name(language)
            style = get_style_by_name(self.config.get("code_style", "monokai"))
            
            tokens = list(lexer.get_tokens(code))
            lines = []
            current_line = []
            for t_type, t_val in tokens:
                if '\n' in t_val:
                    parts = t_val.split('\n')
                    for idx, part in enumerate(parts):
                        if idx > 0:
                            lines.append(current_line)
                            current_line = []
                        if part:
                            current_line.append((t_type, part))
                else:
                    if t_val:
                        current_line.append((t_type, t_val))
            if current_line:
                lines.append(current_line)
            
            if not lines:
                lines = [[]]
                
            # 2. Setup first slide
            slide = self._duplicate_template_slide()
            header_text = f"Question {q_num}:\n{q_statement}"
            header_bottom = self._add_header_to_slide(slide, header_text)
            
            left = Inches(self.config.get("safe_zone_left", 0.5))
            top = header_bottom + Inches(0.2)
            
            max_width = self.prs.slide_width - Inches(self.config.get("safe_zone_left", 0.5) + self.config.get("safe_zone_right", 0.5))
            max_height = self.prs.slide_height - top - Inches(self.config.get("safe_zone_bottom", 1.0))
            
            # Resolve background color and fallback text color
            bg_hex = style.background_color.lstrip('#')
            bg_r = int(bg_hex[0:2], 16)
            bg_g = int(bg_hex[2:4], 16)
            bg_b = int(bg_hex[4:6], 16)
            bg_color_rgb = RGBColor(bg_r, bg_g, bg_b)
            
            # Perceive luminance to ensure text readability
            luminance = bg_r * 0.299 + bg_g * 0.587 + bg_b * 0.114
            if luminance < 128:
                default_text_color = RGBColor(240, 240, 240)
            else:
                default_text_color = RGBColor(30, 30, 30)
                
            # Compute lines that can fit on each slide (standard monospace font size 11pt, line height 1.25)
            font_size_pt = 11
            line_height_pt = font_size_pt * 1.25
            max_height_pt = max_height / 12700
            max_lines_per_slide = max(1, int(max_height_pt / line_height_pt) - 1)
            
            total_lines = len(lines)
            width_format = f"{{:{len(str(total_lines))}d}}  "
            
            slide_idx = 0
            start_line_idx = 0
            
            while start_line_idx < total_lines:
                if slide_idx == 0:
                    current_slide = slide
                    current_top = top
                    current_max_h = max_height
                else:
                    current_slide = self._duplicate_template_slide()
                    h_bottom = self._add_header_to_slide(current_slide, f"")
                    current_top = h_bottom + Inches(0.2)
                    current_max_h = self.prs.slide_height - current_top - Inches(self.config.get("safe_zone_bottom", 1.0))
                
                end_line_idx = min(start_line_idx + max_lines_per_slide, total_lines)
                batch_lines = lines[start_line_idx:end_line_idx]
                
                # Create the code textbox container
                textbox = current_slide.shapes.add_textbox(left, current_top, max_width, current_max_h)
                tf = textbox.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.15)
                tf.margin_right = Inches(0.15)
                tf.margin_top = Inches(0.15)
                tf.margin_bottom = Inches(0.15)
                
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = bg_color_rgb
                
                for line_idx, line_tokens in enumerate(batch_lines):
                    if line_idx > 0:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.line_spacing = 1.2
                    
                    # Prepend line number
                    line_num_run = p.add_run()
                    line_num_run.text = width_format.format(start_line_idx + line_idx + 1)
                    line_num_run.font.name = PPTEngine._resolved_font or "Courier New"
                    line_num_run.font.size = Pt(font_size_pt)
                    line_num_run.font.color.rgb = RGBColor(128, 128, 128)
                    
                    # Add code token runs with syntax highlighting styles
                    for token_type, token_val in line_tokens:
                        run = p.add_run()
                        run.text = token_val
                        run.font.name = PPTEngine._resolved_font or "Courier New"
                        run.font.size = Pt(font_size_pt)
                        
                        style_info = style.style_for_token(token_type)
                        color = style_info.get('color')
                        if color:
                            r = int(color[0:2], 16)
                            g = int(color[2:4], 16)
                            b = int(color[4:6], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
                        else:
                            run.font.color.rgb = default_text_color
                            
                        if style_info.get('bold'):
                            run.font.bold = True
                        if style_info.get('italic'):
                            run.font.italic = True
                
                start_line_idx = end_line_idx
                slide_idx += 1

    def add_screenshot_slide(self, q_num: int, screenshot_path: str):
        """Adds a slide containing the execution screenshot."""
        slide = self._duplicate_template_slide()
        self._add_header_to_slide(slide, f"Question {q_num} (Output):")
        
        left = Inches(self.config.get("safe_zone_left", 0.5))
        top = Inches(self.config.get("safe_zone_top", 1.5)) + Inches(1.0) # Header roughly takes 1 inch
        
        max_width = self.prs.slide_width - Inches(self.config.get("safe_zone_left", 0.5) + self.config.get("safe_zone_right", 0.5))
        max_height = self.prs.slide_height - top - Inches(self.config.get("safe_zone_bottom", 1.0))
        
        with Image.open(screenshot_path) as img:
            img_w, img_h = img.size
            
        aspect = img_h / img_w
        
        if (max_width * aspect) > max_height:
            # Constrain by height
            slide.shapes.add_picture(screenshot_path, left, top, height=max_height)
        else:
            # Constrain by width
            slide.shapes.add_picture(screenshot_path, left, top, width=max_width)
            
    def save(self, output_path: str):
        """Saves the presentation."""
        # Optionally, remove the original template slide if it was just a blank placeholder
        # However, deleting slides in python-pptx is complex and prone to breaking references.
        # It's safer to leave it or ask user to provide a 1-slide template.
        self.prs.save(output_path)
